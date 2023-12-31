from pptx import Presentation
from pptx.util import Inches, Pt
import sys
import os
import json
import pdb
TWO_CONTENT_LAYOUT=3
TITLE_AND_CONTENT_LAYOUT=1
def add_slide(prs, title, content, photo_url):
    slide_layout = prs.slide_layouts[3]  # Use the title and content layout

    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    content_box = slide.placeholders[2]
    

    title_box.text = title
    content_box.text = content
    if os.path.isfile(photo_url):
        picture_box = slide.shapes.add_picture(
            photo_url,
            top=Inches(1.7),
            left=Inches(0.5),
            width=Inches(4.5),
            height=Inches(5),
        )

# Itinerary data
# pdb.set_trace()
itinerary = []
photo_folder = "photos"
output_file_name = "output.pptx"
# for testing
if len(sys.argv) == 1:
    itinerary =  [{
        "name": "樟宜机场",
        "date": "2024-01-22",
        "url": "https://www.changiairport.com",
        "description": "新加坡樟宜機場（英語：Singapore Changi Airport）是新加坡唯一民用机场，也是新加坡最大的国际机场，座落于新加坡東區樟宜，占地13平方公里，距市区约17.2公里。",
        "photo": "changi-airport.jpg"
    }]
elif len(sys.argv) == 2:
    json_file = sys.argv[1]
    with open(json_file, 'r', encoding='utf-8') as file:
        itinerary = json.load(file)
else: 
    json_file = sys.argv[1]
    photo_folder = sys.argv[2]
    with open(json_file, 'r', encoding='utf-8') as file:
        itinerary = json.load(file)

# Create a PowerPoint presentation
presentation = Presentation()

# Add slides for each itinerary item
for item in itinerary:
    photo_url = os.path.join(photo_folder, item["photo"]) 
    add_slide(
        presentation,
        f"{item['name']} ({item['date']})",
        item["description"],
        photo_url,
    )

# Save the presentation to a file
presentation.save(output_file_name)
